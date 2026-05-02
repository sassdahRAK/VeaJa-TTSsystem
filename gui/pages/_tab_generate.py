"""gui/pages/_tab_generate.py — Generate tab for the Dashboard."""

import os
import json

from PyQt6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel, QPushButton,
    QTextEdit, QApplication, QSizePolicy
)
from gui._window_shared import scaled  # noqa: F401
from PyQt6.QtCore import Qt, QThread, QObject, pyqtSignal, QByteArray

from gui.pages._flow_layout import FlowLayout as _FlowLayout
from gui.pages._ai_caller import call_ai, get_api_keys, best_provider

# ── Preview image data ────────────────────────────────────────────────────────
_PREVIEWS_PATH = os.path.join(os.path.dirname(__file__), "..", "..", "assets", "generate_previews.json")

def _load_previews() -> dict:
    try:
        path = os.path.normpath(_PREVIEWS_PATH)
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

_PREVIEWS = _load_previews()


# ── Full-bleed background image widget ───────────────────────────────────────

class _HeroBgWidget(QWidget):
    """
    A widget that paints a QPixmap scaled to cover its entire area (like
    CSS background-size: cover).  Children are laid out normally on top.
    """
    def __init__(self, parent=None):
        super().__init__(parent)
        self._pixmap = None

    def set_pixmap(self, px):
        self._pixmap = px
        self.update()

    def paintEvent(self, event):
        super().paintEvent(event)
        if not self._pixmap or self._pixmap.isNull():
            return
        from PyQt6.QtGui import QPainter
        painter = QPainter(self)
        painter.setRenderHint(QPainter.RenderHint.SmoothPixmapTransform)
        # Scale to cover — fill the widget, crop excess
        w, h = self.width(), self.height()
        px = self._pixmap
        scale = max(w / px.width(), h / px.height())
        new_w = int(px.width()  * scale)
        new_h = int(px.height() * scale)
        x = (w - new_w) // 2
        y = (h - new_h) // 2
        scaled = px.scaled(
            new_w, new_h,
            Qt.AspectRatioMode.IgnoreAspectRatio,
            Qt.TransformationMode.SmoothTransformation,
        )
        painter.drawPixmap(x, y, scaled)
        # Dark gradient overlay so the input card text stays readable
        from PyQt6.QtGui import QLinearGradient, QColor, QBrush
        grad = QLinearGradient(0, 0, 0, h)
        grad.setColorAt(0.0, QColor(0, 0, 0, 30))
        grad.setColorAt(0.6, QColor(0, 0, 0, 80))
        grad.setColorAt(1.0, QColor(0, 0, 0, 200))
        painter.fillRect(0, 0, w, h, QBrush(grad))
        painter.end()


# ── Background worker for fetching preview images ─────────────────────────────

class _ImgSignals(QObject):
    loaded = pyqtSignal(bytes)
    failed = pyqtSignal()

class _ImgFetchThread(QThread):
    def __init__(self, url: str, signals: _ImgSignals):
        super().__init__()
        self._url     = url
        self._signals = signals

    def run(self):
        try:
            import urllib.request
            req = urllib.request.Request(
                self._url,
                headers={"User-Agent": "Mozilla/5.0"}
            )
            with urllib.request.urlopen(req, timeout=8) as resp:
                data = resp.read()
            self._signals.loaded.emit(data)
        except Exception:
            self._signals.failed.emit()


# ── Background worker ─────────────────────────────────────────────────────────

class _GenSignals(QObject):
    finished = pyqtSignal(str)

class _GenThread(QThread):
    def __init__(self, prompt, system, mixin, signals):
        super().__init__()
        self._prompt  = prompt
        self._system  = system
        self._mixin   = mixin
        self._signals = signals

    def run(self):
        result = call_ai("generate", self._prompt, self._mixin, self._system)
        self._signals.finished.emit(result)


class GenerateTabMixin:
    """Generate tab builder and all generation/export logic."""

    _GEN_DESCRIPTIONS = [
        "Create a visual poster layout with headline, body and call-to-action.",
        "Write an AI prompt optimised for image or text generation models.",
        "Produce step-by-step instructions or a how-to guide.",
        "Write a short social-media caption with hashtags.",
        "Rewrite or adjust the tone, length or style of existing text.",
        "Generate a multi-slide presentation outline (title + bullet points per slide).",
        "Generate a video script with scenes, narration and on-screen text.",
        "Generate a self-contained mini web page as raw HTML + CSS.",
        "Generate a full web project scaffold: HTML, CSS, JS, and file structure.",
        "Generate an app development plan: architecture, screens, components, and tech stack.",
        "Generate a network topology plan: devices, subnets, protocols, and config snippets.",
        "Generate a quiz with multiple-choice questions, answers, and explanations.",
    ]

    def _build_generate_tab(self) -> QWidget:
        frame = QWidget()
        frame.setObjectName("tabPage")
        lay = QVBoxLayout(frame)
        lay.setContentsMargins(0, 18, 0, 0)
        lay.setSpacing(10)

        # Sub-mode bar (wraps to 2 rows)
        mode_bar = QWidget()
        mode_bar.setObjectName("subTabBar")
        mb_lay = _FlowLayout(mode_bar, h_spacing=6, v_spacing=6)
        mb_lay.setContentsMargins(0, 4, 0, 4)

        self._gen_modes = ["Poster", "Prompt", "Instruction", "Caption", "Adjust",
                           "Slide", "Video", "HTML", "Web", "App", "Network", "Quiz"]
        self._gen_mode_btns: list[QPushButton] = []
        for i, label in enumerate(self._gen_modes):
            btn = QPushButton(label)
            btn.setObjectName("subTabBtn")
            btn.setCheckable(True)
            btn.setChecked(i == 0)
            btn.setFixedHeight(30)
            btn.setMinimumWidth(btn.fontMetrics().horizontalAdvance(label) + 28)
            btn.clicked.connect(lambda _=False, idx=i: self._switch_gen_mode(idx))
            mb_lay.addWidget(btn)
            self._gen_mode_btns.append(btn)
        lay.addWidget(mode_bar)

        # Input area — transparent so the hero background shows through
        input_card = QWidget()
        input_card.setObjectName("featureBox")
        input_card.setAttribute(Qt.WidgetAttribute.WA_TranslucentBackground, False)
        input_card.setStyleSheet(
            "#featureBox, QWidget#featureBox { background: rgba(20,20,20,0.72); "
            "border-radius: 10px; }"
        )
        ic_lay = QVBoxLayout(input_card)
        ic_lay.setContentsMargins(0, 0, 0, 0)
        ic_lay.setSpacing(6)

        self._gen_mode_desc = QLabel(self._gen_mode_description(0))
        self._gen_mode_desc.setObjectName("featureLabel")
        self._gen_mode_desc.setWordWrap(True)
        ic_lay.addWidget(self._gen_mode_desc)

        self._gen_text_input = QTextEdit()
        self._gen_text_input.setObjectName("featureEdit")
        self._gen_text_input.setPlaceholderText(
            "Describe what you want to generate, or paste text / URL here…"
        )
        self._gen_text_input.setFixedHeight(100)
        ic_lay.addWidget(self._gen_text_input)

        # Attachment + action row — all on the same line
        attach_row = QHBoxLayout()
        attach_row.setSpacing(6)
        attach_row.setContentsMargins(0, 4, 0, 0)

        attach_lbl = QLabel("Attach:")
        attach_lbl.setObjectName("featureLabel")
        attach_row.addWidget(attach_lbl)

        # Attach buttons with inline SVG icons rendered via QPixmap
        self._gen_attach_btns: list[tuple] = []   # (btn, svg_body) for theme refresh
        for svg_body, tip, slot in [
            # Image — landscape photo icon
            (
                '<rect x="3" y="5" width="18" height="14" rx="2" stroke="currentColor" stroke-width="1.6" fill="none"/>'
                '<circle cx="8.5" cy="10.5" r="1.5" fill="currentColor"/>'
                '<path d="M3 16l4-4 3 3 3-4 5 5" stroke="currentColor" stroke-width="1.6" stroke-linejoin="round" fill="none"/>',
                "Image", self._gen_attach_image,
            ),
            # PDF — document with lines icon
            (
                '<path d="M6 2h8l4 4v14a2 2 0 0 1-2 2H6a2 2 0 0 1-2-2V4a2 2 0 0 1 2-2z" stroke="currentColor" stroke-width="1.6" fill="none"/>'
                '<path d="M14 2v4h4" stroke="currentColor" stroke-width="1.6" fill="none"/>'
                '<path d="M8 13h8M8 17h5" stroke="currentColor" stroke-width="1.6" stroke-linecap="round"/>',
                "PDF", self._gen_attach_pdf,
            ),
            # Folder — open folder icon
            (
                '<path d="M3 7a2 2 0 0 1 2-2h4l2 2h8a2 2 0 0 1 2 2v8a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2V7z" stroke="currentColor" stroke-width="1.6" fill="none"/>',
                "Folder", self._gen_attach_folder,
            ),
            # URL — chain link icon
            (
                '<path d="M10 13a5 5 0 0 0 7.54.54l3-3a5 5 0 0 0-7.07-7.07l-1.72 1.71" stroke="currentColor" stroke-width="1.6" stroke-linecap="round" fill="none"/>'
                '<path d="M14 11a5 5 0 0 0-7.54-.54l-3 3a5 5 0 0 0 7.07 7.07l1.71-1.71" stroke="currentColor" stroke-width="1.6" stroke-linecap="round" fill="none"/>',
                "URL", self._gen_attach_url,
            ),
        ]:
            btn = QPushButton(f" {tip}")
            btn.setObjectName("btnOutline")
            btn.setFixedHeight(30)
            btn.setToolTip(tip)
            btn.setIcon(self._gen_svg_icon(svg_body))
            from PyQt6.QtCore import QSize
            btn.setIconSize(QSize(15, 15))
            btn.clicked.connect(slot)
            attach_row.addWidget(btn)
            self._gen_attach_btns.append((btn, svg_body))

        attach_row.addStretch()

        # AI provider indicator
        self._gen_ai_lbl = QLabel("")
        self._gen_ai_lbl.setObjectName("settingsLabel")
        self._gen_ai_lbl.setStyleSheet("font-size: 10px; color: #888;")
        attach_row.addWidget(self._gen_ai_lbl)

        clear_btn = QPushButton("Clear")
        clear_btn.setObjectName("btnOutline")
        clear_btn.setFixedSize(72, 30)
        clear_btn.clicked.connect(self._clear_generate)
        attach_row.addWidget(clear_btn)

        self._gen_btn = QPushButton("Generate")
        self._gen_btn.setObjectName("btnPrimary")
        self._gen_btn.setFixedSize(100, 30)
        self._gen_btn.clicked.connect(self._run_generate)
        attach_row.addWidget(self._gen_btn)

        ic_lay.addLayout(attach_row)

        self._gen_attachments: list[str] = []
        self._gen_attach_lbl = QLabel("")
        self._gen_attach_lbl.setObjectName("settingsLabel")
        self._gen_attach_lbl.setWordWrap(True)
        self._gen_attach_lbl.setStyleSheet("font-size: 11px;")
        self._gen_attach_lbl.setVisible(False)
        ic_lay.addWidget(self._gen_attach_lbl)

        # ── Hero container: full-bleed background image + input card overlay ──
        # The hero widget fills all remaining vertical space. The background
        # image is painted to cover the entire area; the input card sits at
        # the bottom, semi-transparent, so the image shows behind it.
        self._gen_hero = _HeroBgWidget()
        self._gen_hero.setSizePolicy(
            QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding
        )
        hero_lay = QVBoxLayout(self._gen_hero)
        hero_lay.setContentsMargins(0, 0, 0, 0)
        hero_lay.setSpacing(0)
        hero_lay.addStretch(1)          # push input card to the bottom
        hero_lay.addWidget(input_card)  # input card sits at the bottom of the hero

        lay.addWidget(self._gen_hero, 1)

        self._gen_img_threads: list = []   # keep thread refs alive
        # Load the first mode's preview immediately
        self._load_gen_preview(0)

        # Output
        self._gen_out_lbl = QLabel("Output")
        self._gen_out_lbl.setObjectName("featureLabel")
        self._gen_out_lbl.setVisible(False)
        lay.addWidget(self._gen_out_lbl)

        self._gen_output = QTextEdit()
        self._gen_output.setObjectName("featureEditReadOnly")
        self._gen_output.setReadOnly(True)
        self._gen_output.setPlaceholderText("Generated content will appear here…")
        self._gen_output.setVisible(False)
        lay.addWidget(self._gen_output, 1)

        # Export row (wraps to 2 rows)
        self._gen_export_row = QWidget()
        self._gen_export_row.setVisible(False)
        er_lay = _FlowLayout(self._gen_export_row, h_spacing=6, v_spacing=6)
        er_lay.setContentsMargins(0, 4, 0, 4)

        er_lay.addWidget(QLabel("Export as:", objectName="featureLabel"))

        for label, icon, slot in [
            ("Copy",    "📋", self._gen_export_copy),
            ("PDF",     "📄", self._gen_export_pdf),
            ("DOCX",    "📝", self._gen_export_docx),
            ("JPG",     "🖼", self._gen_export_jpg),
            ("PNG",     "🖼", self._gen_export_png),
            ("PPTX",    "📊", self._gen_export_pptx),
            ("HTML",    "🌐", self._gen_export_html),
            ("QR Code", "⬛", self._gen_export_qr),
        ]:
            btn = QPushButton(f"{icon}  {label}")
            btn.setObjectName("btnOutline")
            btn.setFixedHeight(30)
            btn.setMinimumWidth(btn.fontMetrics().horizontalAdvance(f"{icon}  {label}") + 28)
            btn.clicked.connect(slot)
            er_lay.addWidget(btn)

        self._gen_preview_btn = QPushButton("🔍  Preview in browser")
        self._gen_preview_btn.setObjectName("btnPrimary")
        self._gen_preview_btn.setFixedHeight(30)
        self._gen_preview_btn.setVisible(False)
        self._gen_preview_btn.clicked.connect(self._gen_preview_html)
        er_lay.addWidget(self._gen_preview_btn)
        lay.addWidget(self._gen_export_row)
        return frame

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _gen_svg_icon(self, svg_body: str, size: int = 16):
        """Render an SVG path string into a QIcon, respecting dark/light theme."""
        from PyQt6.QtSvg import QSvgRenderer
        from PyQt6.QtGui import QPixmap, QPainter, QIcon, QColor
        from PyQt6.QtWidgets import QApplication
        is_dark = getattr(self, "_dark", True)
        color = "#ffffff" if is_dark else "#333333"
        # Replace currentColor with the resolved color so it renders correctly
        svg_colored = svg_body.replace("currentColor", color)
        svg = (
            f'<svg viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg">'
            f'<g stroke="{color}" fill="none">{svg_colored}</g>'
            f'</svg>'
        ).encode()
        app = QApplication.instance()
        dpr = app.primaryScreen().devicePixelRatio() if app and app.primaryScreen() else 1.0
        phys = int(size * dpr)
        px = QPixmap(phys, phys)
        px.fill(QColor(0, 0, 0, 0))
        renderer = QSvgRenderer(svg)
        p = QPainter(px)
        p.setRenderHint(QPainter.RenderHint.Antialiasing)
        renderer.render(p)
        p.end()
        px.setDevicePixelRatio(dpr)
        return QIcon(px)

    def _gen_mode_description(self, idx: int) -> str:
        return self._GEN_DESCRIPTIONS[idx] if idx < len(self._GEN_DESCRIPTIONS) else ""

    def refresh_generate_icons(self):
        """Re-render theme-sensitive attach icons (called on theme toggle)."""
        if not hasattr(self, "_gen_attach_btns"):
            return
        from PyQt6.QtCore import QSize
        for btn, svg_body in self._gen_attach_btns:
            btn.setIcon(self._gen_svg_icon(svg_body))
            btn.setIconSize(QSize(15, 15))

    def _switch_gen_mode(self, idx: int):
        for i, btn in enumerate(self._gen_mode_btns):
            btn.setChecked(i == idx)
        self._gen_mode_desc.setText(self._gen_mode_description(idx))
        # Reload preview image for the new mode (only if output is hidden)
        if hasattr(self, "_gen_hero") and self._gen_hero.isVisible():
            self._load_gen_preview(idx)

    def _current_gen_mode(self) -> str:
        for i, btn in enumerate(self._gen_mode_btns):
            if btn.isChecked():
                return self._gen_modes[i]
        return "Poster"

    def _update_attach_label(self):
        if self._gen_attachments:
            self._gen_attach_lbl.setText("Attached: " + "  ·  ".join(self._gen_attachments))
            self._gen_attach_lbl.setVisible(True)
        else:
            self._gen_attach_lbl.setVisible(False)

    def _load_gen_preview(self, idx: int):
        """Fetch and display the preview image for the given mode index."""
        if not hasattr(self, "_gen_hero"):
            return
        mode = self._gen_modes[idx] if idx < len(self._gen_modes) else "Poster"
        entry = _PREVIEWS.get(mode, {})
        url = entry.get("url", "")
        if not url:
            self._gen_hero.set_pixmap(None)
            return
        signals = _ImgSignals()
        signals.loaded.connect(lambda data, m=mode: self._on_preview_loaded(data, m))
        signals.failed.connect(lambda m=mode: self._on_preview_failed(m))
        t = _ImgFetchThread(url, signals)
        self._gen_img_threads.append(t)
        t.start()

    def _on_preview_loaded(self, data: bytes, mode: str):
        """Paint the downloaded image as the hero background."""
        if not hasattr(self, "_gen_hero"):
            return
        if self._current_gen_mode() != mode:
            return
        from PyQt6.QtGui import QPixmap
        px = QPixmap()
        px.loadFromData(QByteArray(data))
        if not px.isNull():
            self._gen_hero.set_pixmap(px)

    def _on_preview_failed(self, mode: str):
        if not hasattr(self, "_gen_hero"):
            return
        if self._current_gen_mode() != mode:
            return
        self._gen_hero.set_pixmap(None)

    def _gen_attach_image(self):
        from PyQt6.QtWidgets import QFileDialog
        path, _ = QFileDialog.getOpenFileName(None, "Attach image", "",
                                              "Images (*.jpg *.jpeg *.png *.webp *.bmp *.gif)")
        if path:
            self._gen_attachments.append(path)
            self._update_attach_label()

    def _gen_attach_pdf(self):
        from PyQt6.QtWidgets import QFileDialog
        path, _ = QFileDialog.getOpenFileName(None, "Attach PDF", "", "PDF (*.pdf)")
        if path:
            self._gen_attachments.append(path)
            self._update_attach_label()

    def _gen_attach_folder(self):
        from PyQt6.QtWidgets import QFileDialog
        path = QFileDialog.getExistingDirectory(None, "Attach folder")
        if path:
            self._gen_attachments.append(path + "/")
            self._update_attach_label()

    def _gen_attach_url(self):
        from PyQt6.QtWidgets import QInputDialog
        url, ok = QInputDialog.getText(None, "Attach URL", "Enter URL:")
        if ok and url.strip():
            self._gen_attachments.append(url.strip())
            self._update_attach_label()

    def _clear_generate(self):
        self._gen_text_input.clear()
        self._gen_attachments.clear()
        self._update_attach_label()
        self._gen_output.clear()
        self._gen_output.setVisible(False)
        self._gen_out_lbl.setVisible(False)
        self._gen_export_row.setVisible(False)
        # Show hero preview again after clearing
        if hasattr(self, "_gen_hero"):
            self._gen_hero.setVisible(True)
            idx = next((i for i, b in enumerate(self._gen_mode_btns) if b.isChecked()), 0)
            self._load_gen_preview(idx)

    def _run_generate(self):
        text = self._gen_text_input.toPlainText().strip()
        mode = self._current_gen_mode()
        if not text and not self._gen_attachments:
            return

        # Check if AI key available
        keys = get_api_keys(self)
        provider_result = best_provider("generate", keys)

        if provider_result:
            # Use AI
            provider, _ = provider_result
            self._gen_ai_lbl.setText(f"✦ {provider.title()}")
            self._gen_btn.setEnabled(False)
            self._gen_btn.setText("Generating…")
            self._gen_output.setPlainText("Generating…")
            self._gen_output.setVisible(True)
            self._gen_out_lbl.setVisible(True)
            # Hide hero while output is shown
            if hasattr(self, "_gen_hero"):
                self._gen_hero.setVisible(False)

            system, prompt = self._build_gen_prompt(mode, text)
            signals = _GenSignals()
            signals.finished.connect(self._on_gen_finished)
            self._gen_thread = _GenThread(prompt, system, self, signals)
            self._gen_thread.start()
        else:
            # Fallback template
            self._gen_ai_lbl.setText("No AI key — using template")
            result = self._generate_content(mode, text, self._gen_attachments)
            self._gen_output.setPlainText(result)
            self._gen_output.setVisible(True)
            self._gen_out_lbl.setVisible(True)
            self._gen_export_row.setVisible(True)
            if hasattr(self, "_gen_preview_btn"):
                self._gen_preview_btn.setVisible(mode == "HTML")
            if hasattr(self, "_gen_hero"):
                self._gen_hero.setVisible(False)

    def _on_gen_finished(self, result: str):
        self._gen_btn.setEnabled(True)
        self._gen_btn.setText("Generate")
        self._gen_output.setPlainText(result)
        self._gen_export_row.setVisible(True)
        if hasattr(self, "_gen_preview_btn"):
            self._gen_preview_btn.setVisible(self._current_gen_mode() == "HTML")

    def _build_gen_prompt(self, mode: str, text: str) -> tuple[str, str]:
        """Return (system, user_prompt) for the given generate mode."""
        attach_note = ""
        if self._gen_attachments:
            attach_note = "\n\nAttached sources:\n" + "\n".join(f"  • {a}" for a in self._gen_attachments)

        systems = {
            "Poster":      "You are a professional copywriter. Create a visual poster layout with a compelling headline, body text, and call-to-action. Use clear sections.",
            "Prompt":      "You are an expert prompt engineer. Write an optimised AI prompt for the given task. Include role, context, task, format, and constraints.",
            "Instruction": "You are a technical writer. Write clear, numbered step-by-step instructions. Include prerequisites, steps, and tips.",
            "Caption":     "You are a social media expert. Write an engaging caption with relevant hashtags. Keep it punchy and on-brand.",
            "Adjust":      "You are an expert editor. Rewrite the given text to improve clarity, tone, and style. Preserve the original meaning.",
            "Slide":       "You are a presentation designer. Create a detailed slide-by-slide outline with titles and bullet points for each slide.",
            "Video":       "You are a video scriptwriter. Write a complete video script with scenes, timestamps, narration, and on-screen text.",
            "HTML":        "You are a web developer. Generate a complete, self-contained HTML page with embedded CSS. Make it modern and responsive.",
            "Web":         "You are a senior web architect. Generate a complete web project scaffold with file structure, tech stack, and starter code for each file.",
            "App":         "You are a mobile/web app architect. Generate a detailed app development plan with architecture, screens, components, and tech stack.",
            "Network":     "You are a network engineer. Generate a detailed network topology plan with devices, subnets, protocols, and config snippets.",
            "Quiz":        "You are an educator. Generate a 5-question quiz with multiple choice answers, correct answers marked, and explanations.",
        }
        system = systems.get(mode, "You are a helpful assistant.")
        prompt = f"{text}{attach_note}"
        return system, prompt

    def _generate_content(self, mode: str, text: str, attachments: list) -> str:
        attach_note = (
            "\n\nAttached sources:\n" + "\n".join(f"  • {a}" for a in attachments)
            if attachments else ""
        )
        templates = {
            "Poster": (
                f"━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
                f"  HEADLINE\n  {text[:60] or 'Your headline here'}\n\n"
                f"  BODY\n  {text or 'Describe your message here.'}\n\n"
                f"  CALL TO ACTION\n  Learn more → yourlink.com\n"
                f"━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"
            ),
            "Prompt": (
                f"You are a helpful assistant.\n\nTask: {text or 'Describe the task.'}\n\n"
                f"Requirements:\n- Be concise and clear\n- Use structured output\n"
                f"- Tone: professional\n\nOutput format: [specify format here]"
            ),
            "Instruction": (
                f"# How to: {text[:50] or 'Your topic'}\n\n"
                f"## Prerequisites\n- Item 1\n- Item 2\n\n"
                f"## Steps\n1. First step — {text[:80] or 'describe step 1'}\n"
                f"2. Second step\n3. Third step\n\n## Notes\n- Tip or warning here."
            ),
            "Caption": f"{text[:120] or 'Your caption here.'} ✨\n\n#YourBrand #Topic #Trending",
            "Adjust": (
                f"[Adjusted version of your text]\n\n"
                f"{text or 'Paste text above to adjust its tone, length or style.'}\n\n"
                f"— Connect an AI API key in My API Key to enable live adjustment."
            ),
            "Slide": self._gen_slide_template(text),
            "Video": self._gen_video_template(text),
            "HTML":  self._gen_html_template(text),
            "Web":     self._gen_web_template(text),
            "App":     self._gen_app_template(text),
            "Network": self._gen_network_template(text),
            "Quiz":    self._gen_quiz_template(text),
        }
        return templates.get(mode, text) + attach_note

    # ── Export ────────────────────────────────────────────────────────────────

    def _gen_export_copy(self):
        text = self._gen_output.toPlainText()
        if text:
            QApplication.clipboard().setText(text)

    def _gen_export_pdf(self):   self._gen_save_file("PDF (*.pdf)", ".pdf", self._write_pdf)
    def _gen_export_docx(self):  self._gen_save_file("Word Document (*.docx)", ".docx", self._write_docx)
    def _gen_export_jpg(self):   self._gen_save_file("JPEG Image (*.jpg)", ".jpg", self._write_image_jpg)
    def _gen_export_png(self):   self._gen_save_file("PNG Image (*.png)", ".png", self._write_image_png)
    def _gen_export_pptx(self):  self._gen_save_file("PowerPoint (*.pptx)", ".pptx", self._write_pptx)
    def _gen_export_qr(self):    self._gen_save_file("PNG Image (*.png)", ".png", self._write_qr)
    def _gen_export_html(self):  self._gen_save_file("HTML File (*.html)", ".html", self._write_html)

    def _gen_save_file(self, filter_str: str, ext: str, writer):
        from PyQt6.QtWidgets import QFileDialog, QMessageBox
        path, _ = QFileDialog.getSaveFileName(None, "Export", f"generated{ext}", filter_str)
        if path:
            try:
                writer(path)
            except Exception as e:
                QMessageBox.warning(None, "Export failed", str(e))

    def _write_pdf(self, path: str):
        from PyQt6.QtGui import QTextDocument
        from PyQt6.QtPrintSupport import QPrinter
        printer = QPrinter(QPrinter.PrinterMode.HighResolution)
        printer.setOutputFormat(QPrinter.OutputFormat.PdfFormat)
        printer.setOutputFileName(path)
        doc = QTextDocument()
        doc.setPlainText(self._gen_output.toPlainText())
        doc.print(printer)

    def _write_docx(self, path: str):
        try:
            from docx import Document
            doc = Document()
            for line in self._gen_output.toPlainText().splitlines():
                doc.add_paragraph(line)
            doc.save(path)
        except ImportError:
            with open(path, "w", encoding="utf-8") as f:
                f.write(self._gen_output.toPlainText())

    def _write_image_jpg(self, path: str): self._render_output_image(path, "JPEG")
    def _write_image_png(self, path: str): self._render_output_image(path, "PNG")

    def _render_output_image(self, path: str, fmt: str):
        from PyQt6.QtGui import QPixmap, QPainter, QFont, QColor
        from PyQt6.QtCore import QRect
        lines = self._gen_output.toPlainText().splitlines()
        W, line_h, pad = 800, 22, 30
        H = max(pad * 2 + len(lines) * line_h + 20, 200)
        px = QPixmap(W, H)
        px.fill(QColor("#ffffff"))
        p = QPainter(px)
        p.setFont(QFont("Segoe UI", 10))
        p.setPen(QColor("#1a1a1a"))
        y = pad
        for line in lines:
            p.drawText(QRect(pad, y, W - pad * 2, line_h), Qt.AlignmentFlag.AlignLeft, line)
            y += line_h
        p.end()
        px.save(path, fmt)

    def _write_pptx(self, path: str):
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = self._current_gen_mode()
            slide.placeholders[1].text_frame.text = self._gen_output.toPlainText()[:500]
            prs.save(path)
        except ImportError:
            with open(path, "w", encoding="utf-8") as f:
                f.write(self._gen_output.toPlainText())

    def _write_qr(self, path: str):
        try:
            import qrcode
            qrcode.make(self._gen_output.toPlainText()[:300]).save(path)
        except ImportError:
            self._render_output_image(path, "PNG")

    def _write_html(self, path: str):
        with open(path, "w", encoding="utf-8") as f:
            f.write(self._gen_output.toPlainText())

    def _gen_preview_html(self):
        import tempfile
        from PyQt6.QtGui import QDesktopServices
        from PyQt6.QtCore import QUrl
        tmp = tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8")
        tmp.write(self._gen_output.toPlainText())
        tmp.close()
        QDesktopServices.openUrl(QUrl.fromLocalFile(tmp.name))

    # ── Templates ─────────────────────────────────────────────────────────────

    def _gen_slide_template(self, text: str) -> str:
        title = text[:60] or "Presentation Title"
        return (
            f"SLIDE 1 — Title\n  Title:    {title}\n  Subtitle: Your subtitle here\n\n"
            f"SLIDE 2 — Introduction\n  • What is this about?\n  • Why does it matter?\n\n"
            f"SLIDE 3 — Key Point 1\n  • {text[:80] or 'First key point'}\n  • Supporting detail\n\n"
            f"SLIDE 4 — Key Point 2\n  • Second key point\n  • Supporting detail\n\n"
            f"SLIDE 5 — Key Point 3\n  • Third key point\n  • Supporting detail\n\n"
            f"SLIDE 6 — Summary\n  • Recap of main points\n  • Call to action\n\n"
            f"SLIDE 7 — Thank You\n  Contact: your@email.com\n\n"
            f"— Connect an AI API key to generate real slide content."
        )

    def _gen_video_template(self, text: str) -> str:
        title = text[:50] or "Video Title"
        return (
            f"VIDEO SCRIPT — {title}\n{'─'*50}\n\n"
            f"[SCENE 1 — Hook  |  0:00–0:10]\n"
            f"  ON SCREEN:  \"{title}\"\n"
            f"  NARRATION:  \"Did you know that {text[:60] or '...'}?\"\n"
            f"  MUSIC:      Upbeat intro\n\n"
            f"[SCENE 2 — Problem  |  0:10–0:30]\n"
            f"  ON SCREEN:  Problem statement\n  NARRATION:  \"Many people struggle with...\"\n\n"
            f"[SCENE 3 — Solution  |  0:30–1:00]\n"
            f"  ON SCREEN:  Step-by-step breakdown\n  NARRATION:  \"Here's how to solve it...\"\n\n"
            f"[SCENE 4 — Demo  |  1:00–1:30]\n"
            f"  ON SCREEN:  Screen recording\n  NARRATION:  \"Let me show you exactly how...\"\n\n"
            f"[SCENE 5 — CTA  |  1:30–1:45]\n"
            f"  ON SCREEN:  Subscribe button\n  NARRATION:  \"Like and subscribe!\"\n\n"
            f"— Connect an AI API key to generate a real video script."
        )

    def _gen_html_template(self, text: str) -> str:
        title = text[:50] or "My Page"
        body  = text or "Add your content here."
        return (
            "<!DOCTYPE html>\n<html lang=\"en\">\n<head>\n"
            "  <meta charset=\"UTF-8\">\n"
            "  <meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\">\n"
            f"  <title>{title}</title>\n"
            "  <style>\n"
            "    *, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }\n"
            "    body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;\n"
            "           background: #f8f9fa; color: #1a1a1a; line-height: 1.7; }\n"
            "    header { background: #1a1a1a; color: #fff; padding: 2rem; text-align: center; }\n"
            "    header h1 { font-size: 2rem; font-weight: 700; }\n"
            "    main { max-width: 800px; margin: 3rem auto; padding: 0 1.5rem; }\n"
            "    section { background: #fff; border-radius: 12px; padding: 2rem;\n"
            "              margin-bottom: 1.5rem; box-shadow: 0 2px 8px rgba(0,0,0,.06); }\n"
            "    h2 { font-size: 1.3rem; margin-bottom: 1rem; }\n"
            "    .cta { display: inline-block; margin-top: 1.5rem; background: #1a1a1a;\n"
            "           color: #fff; padding: .75rem 2rem; border-radius: 8px;\n"
            "           text-decoration: none; font-weight: 600; }\n"
            "    footer { text-align: center; padding: 2rem; color: #999; font-size: .85rem; }\n"
            "  </style>\n</head>\n<body>\n"
            f"  <header><h1>{title}</h1><p>Generated by Veaja</p></header>\n"
            "  <main>\n"
            f"    <section><h2>About</h2><p>{body}</p>"
            "<a class=\"cta\" href=\"#\">Get Started</a></section>\n"
            "    <section><h2>Features</h2>"
            "<p>• Feature one<br>• Feature two<br>• Feature three</p></section>\n"
            "    <section><h2>Contact</h2><p>your@email.com</p></section>\n"
            "  </main>\n"
            f"  <footer>© 2025 {title}. All rights reserved.</footer>\n"
            "</body>\n</html>"
        )

    def _gen_web_template(self, text: str) -> str:
        title = text[:50] or "My Web Project"
        return (
            f"# Web Project: {title}\n"
            f"{'─'*50}\n\n"
            f"## Project Structure\n"
            f"  {title.lower().replace(' ', '-')}/\n"
            f"  ├── index.html          # Main entry point\n"
            f"  ├── about.html          # About page\n"
            f"  ├── css/\n"
            f"  │   ├── style.css       # Global styles\n"
            f"  │   └── responsive.css  # Mobile breakpoints\n"
            f"  ├── js/\n"
            f"  │   ├── main.js         # App logic\n"
            f"  │   └── api.js          # API calls\n"
            f"  ├── assets/\n"
            f"  │   ├── images/\n"
            f"  │   └── fonts/\n"
            f"  └── README.md\n\n"
            f"## Tech Stack\n"
            f"  Frontend:  HTML5 · CSS3 · Vanilla JS (or React/Vue)\n"
            f"  Styling:   Tailwind CSS / Bootstrap\n"
            f"  Build:     Vite / Webpack\n"
            f"  Deploy:    Vercel / Netlify / GitHub Pages\n\n"
            f"## Key Pages\n"
            f"  1. Home       — Hero, features, CTA\n"
            f"  2. About      — Team, mission, story\n"
            f"  3. Services   — {text[:60] or 'List your services'}\n"
            f"  4. Contact    — Form, map, social links\n\n"
            f"## SEO & Performance\n"
            f"  • Meta tags, Open Graph, sitemap.xml\n"
            f"  • Lazy loading images\n"
            f"  • Minified CSS/JS in production\n\n"
            f"— Connect an AI API key to generate real code for each file."
        )

    def _gen_app_template(self, text: str) -> str:
        title = text[:50] or "My App"
        return (
            f"# App Development Plan: {title}\n"
            f"{'─'*50}\n\n"
            f"## Overview\n"
            f"  {text or 'Describe your app idea here.'}\n\n"
            f"## Platform\n"
            f"  [ ] Mobile (iOS + Android) — React Native / Flutter\n"
            f"  [ ] Desktop — Electron / PyQt / Tauri\n"
            f"  [ ] Web App — React / Vue / Next.js\n\n"
            f"## Architecture\n"
            f"  Frontend:  Component-based UI\n"
            f"  Backend:   REST API / GraphQL\n"
            f"  Database:  PostgreSQL / SQLite / Firebase\n"
            f"  Auth:      JWT / OAuth2\n"
            f"  Storage:   S3 / Cloudinary\n\n"
            f"## Screens / Views\n"
            f"  1. Splash / Onboarding\n"
            f"  2. Login / Register\n"
            f"  3. Dashboard / Home\n"
            f"  4. {text[:40] or 'Main feature screen'}\n"
            f"  5. Settings / Profile\n\n"
            f"## Key Components\n"
            f"  • Navigation (bottom tab / drawer)\n"
            f"  • Data fetching (React Query / SWR)\n"
            f"  • State management (Redux / Zustand / Provider)\n"
            f"  • Push notifications\n"
            f"  • Offline support\n\n"
            f"## Development Phases\n"
            f"  Phase 1 — MVP (4 weeks): Core screens + auth\n"
            f"  Phase 2 — Beta (3 weeks): API integration + testing\n"
            f"  Phase 3 — Launch (2 weeks): Polish + store submission\n\n"
            f"— Connect an AI API key to generate real code scaffolding."
        )

    def _gen_network_template(self, text: str) -> str:
        title = text[:50] or "Network Design"
        return (
            f"# Network Topology Plan: {title}\n"
            f"{'─'*50}\n\n"
            f"## Overview\n"
            f"  {text or 'Describe your network requirements here.'}\n\n"
            f"## Network Diagram (Text)\n"
            f"  Internet\n"
            f"      │\n"
            f"  [Firewall / Router]  — WAN: 203.0.113.1/30\n"
            f"      │\n"
            f"  [Core Switch L3]     — VLAN trunk\n"
            f"    ├── VLAN 10 (LAN)     192.168.10.0/24\n"
            f"    ├── VLAN 20 (DMZ)     192.168.20.0/24\n"
            f"    ├── VLAN 30 (Mgmt)    192.168.30.0/24\n"
            f"    └── VLAN 40 (WiFi)    192.168.40.0/24\n\n"
            f"## Devices\n"
            f"  • Firewall:      Cisco ASA / pfSense / FortiGate\n"
            f"  • Core Switch:   Cisco Catalyst 3850 / Juniper EX\n"
            f"  • Access Points: Ubiquiti UniFi / Cisco Meraki\n"
            f"  • Servers:       Web, DB, DNS, NTP\n\n"
            f"## Protocols & Services\n"
            f"  Routing:   OSPF / BGP (external)\n"
            f"  DHCP:      Scope per VLAN\n"
            f"  DNS:       Internal + forwarder to 8.8.8.8\n"
            f"  VPN:       IPSec / WireGuard for remote access\n"
            f"  Monitoring: SNMP / Zabbix / Grafana\n\n"
            f"## Security\n"
            f"  • ACLs on inter-VLAN routing\n"
            f"  • IDS/IPS on DMZ\n"
            f"  • 802.1X port authentication\n"
            f"  • NTP sync for log correlation\n\n"
            f"## Sample Config Snippet (Cisco IOS)\n"
            f"  interface Vlan10\n"
            f"   ip address 192.168.10.1 255.255.255.0\n"
            f"   ip helper-address 192.168.30.10\n"
            f"   no shutdown\n\n"
            f"— Connect an AI API key to generate full device configs."
        )

    def _gen_quiz_template(self, text: str) -> str:
        topic = text[:60] or "General Knowledge"
        return (
            f"# Quiz: {topic}\n"
            f"{'─'*50}\n\n"
            f"Instructions: Choose the best answer for each question.\n\n"
            f"──────────────────────────────────────\n"
            f"Q1. What is the main concept of '{topic}'?\n\n"
            f"   A) First option\n"
            f"   B) Second option\n"
            f"   C) Third option  ✓\n"
            f"   D) Fourth option\n\n"
            f"   ✎ Explanation: The correct answer is C because...\n\n"
            f"──────────────────────────────────────\n"
            f"Q2. Which of the following best describes {topic[:30] or 'this topic'}?\n\n"
            f"   A) Incorrect description\n"
            f"   B) Correct description  ✓\n"
            f"   C) Partially correct\n"
            f"   D) Unrelated answer\n\n"
            f"   ✎ Explanation: B is correct because...\n\n"
            f"──────────────────────────────────────\n"
            f"Q3. True or False: {text[:80] or 'This statement is about the topic.'}\n\n"
            f"   A) True  ✓\n"
            f"   B) False\n\n"
            f"   ✎ Explanation: This is true because...\n\n"
            f"──────────────────────────────────────\n"
            f"Q4. Fill in the blank: _______ is a key part of {topic[:30] or 'this subject'}.\n\n"
            f"   Answer: [key concept here]\n\n"
            f"   ✎ Explanation: ...\n\n"
            f"──────────────────────────────────────\n"
            f"Q5. Short answer: Explain {topic[:40] or 'the main idea'} in your own words.\n\n"
            f"   Model answer: ...\n\n"
            f"{'─'*50}\n"
            f"Score: __ / 5\n\n"
            f"— Connect an AI API key to generate a real quiz on any topic."
        )
